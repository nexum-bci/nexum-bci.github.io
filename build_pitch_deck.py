#!/usr/bin/env python3
"""
Nexum One Investor Pitch Deck — Professional 14-Slide Rebuild
逐念而行 · 神经重连AI系统
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import os

# ═══════════════════════════════════════════════════
# DESIGN TOKENS
# ═══════════════════════════════════════════════════
BG_DARK     = RGBColor(0x0A, 0x0F, 0x1A)
CARD_DARK   = RGBColor(0x15, 0x1D, 0x2E)
ACCENT      = RGBColor(0x00, 0x96, 0xF0)
ACCENT_DARK = RGBColor(0x00, 0x56, 0x90)
ACCENT_LIGHT = RGBColor(0x7D, 0xD3, 0xFC)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DARK   = RGBColor(0x64, 0x74, 0x8B)
GRAY_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_CARD   = RGBColor(0xF8, 0xFA, 0xFC)
HIGHLIGHT   = RGBColor(0xE8, 0xF4, 0xFD)
DARK_TEXT    = RGBColor(0x0A, 0x0F, 0x1A)
ROW_ALT     = RGBColor(0xF8, 0xFA, 0xFC)
DARK_CARD   = RGBColor(0x1A, 0x22, 0x34)
TEXT_DARK    = RGBColor(0x1E, 0x29, 0x3B)

# Slide dimensions (16:9 widescreen)
W = Inches(13.333)
H = Inches(7.5)
M = Inches(0.8)       # consistent margin
CW = W - 2 * M        # content width

# Base output directory
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
DOCS_DIR = os.path.join(SCRIPT_DIR, 'docs')
OUTPUT = os.path.join(DOCS_DIR, 'Nexum_Pitch_Deck.pptx')

# ═══════════════════════════════════════════════════
# UTILITY FUNCTIONS
# ═══════════════════════════════════════════════════

def tb(slide, left, top, width, height):
    """Add a textbox and return its text frame."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=Pt(13), color=GRAY, bold=False, align=PP_ALIGN.LEFT,
         space_after=Pt(6), space_before=Pt(0), first=False):
    """Add a paragraph to a text frame."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    p.space_after = space_after
    p.space_before = space_before
    return p


def text(slide, left, top, width, height, txt, size=Pt(13), color=GRAY,
         bold=False, align=PP_ALIGN.LEFT):
    """Single-paragraph text box."""
    tf = tb(slide, left, top, width, height)
    para(tf, txt, size, color, bold, align, first=True)
    return tf


def rect(slide, left, top, width, height, color, radius=None):
    """Add a filled rectangle, optionally rounded."""
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape


def line(slide, left, top, width, color, thickness=Pt(2)):
    """Add a thin horizontal bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def gradient_bg(slide, c1, c2, angle=90.0):
    """Fill entire slide with a gradient."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.line.fill.background()
    f = shape.fill
    f.gradient()
    f.gradient_angle = angle
    f.gradient_stops[0].color.rgb = c1
    f.gradient_stops[1].color.rgb = c2
    return shape


def card(slide, left, top, width, height, bg=WHITE, border_color=None):
    """Card with optional border."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def accent_bar(slide, top, width, color=ACCENT, thickness=Pt(6)):
    """Full-width accent bar."""
    line(slide, 0, top, width, color, thickness)


def content_header(slide, num, title, y_num=Inches(0.38), y_title=Inches(0.72),
                   y_line=Inches(1.22)):
    """Standard content slide header: section number + title + divider line."""
    text(slide, M, y_num, Inches(1.5), Inches(0.3), num, Pt(14), ACCENT, bold=True)
    text(slide, M, y_title, CW, Inches(0.55), title, Pt(30), DARK_TEXT, bold=True)
    line(slide, M, y_line, Inches(1.8), ACCENT, Pt(3))


def make_content_slide(prs, num, title):
    """Create a content slide with white bg, accent bar, and header."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, WHITE)
    accent_bar(s, 0, W, ACCENT, Pt(6))
    content_header(s, num, title)
    return s


def circle_text(slide, left, top, diameter, text_content, bg=ACCENT, txt_color=WHITE):
    """Add a circle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, diameter, diameter)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text_content
    p.font.size = Pt(int(diameter / Inches(1) * 14))
    p.font.color.rgb = txt_color
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_table(slide, left, top, width, height, headers, rows, col_widths,
              font_size=Pt(9)):
    """Professional table with blue header and alternating rows."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = tbl_shape.table
    for i, cw in enumerate(col_widths):
        tbl.columns[i].width = cw

    # Header row
    for j, hdr in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = hdr
        for p in cell.text_frame.paragraphs:
            p.font.size = font_size
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = 'Arial'
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = font_size
                p.font.color.rgb = DARK_TEXT if j == 1 else GRAY_DARK
                p.font.name = 'Arial'
            if i % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = ROW_ALT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return tbl


def metric_card(slide, left, top, width, height, value_text, label_text):
    """Metric display: large value + label."""
    card(slide, left, top, width, height, HIGHLIGHT, ACCENT)
    text(slide, left + Inches(0.1), top + Inches(0.12), width - Inches(0.2),
         Inches(0.55), value_text, Pt(26), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.1), top + Inches(0.72), width - Inches(0.2),
         Inches(0.35), label_text, Pt(10), DARK_TEXT, align=PP_ALIGN.CENTER)


def try_img(slide, path, left, top, width, height=None):
    """Add image if it exists; return True on success."""
    h = height or width * 0.75
    if os.path.exists(path) and os.path.getsize(path) > 1000:
        slide.shapes.add_picture(path, left, top, width, h)
        return True
    return False


# ═══════════════════════════════════════════════════
# HELPERS: Rich paragraph
# ═══════════════════════════════════════════════════

def rich_textbox(slide, left, top, width, height, segments):
    """
    segments = [(text, size, color, bold), ...]
    Returns the text frame.
    """
    tf = tb(slide, left, top, width, height)
    for i, (txt, sz, clr, bd) in enumerate(segments):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = txt
            p.font.size = sz
            p.font.color.rgb = clr
            p.font.bold = bd
            p.font.name = 'Arial'
        else:
            run = tf.paragraphs[0].add_run()
            run.text = txt
            run.font.size = sz
            run.font.color.rgb = clr
            run.font.bold = bd
            run.font.name = 'Arial'
    return tf


# ═══════════════════════════════════════════════════
# INIT PRESENTATION
# ═══════════════════════════════════════════════════

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════
# SLIDE 1 — TITLE (Dark Hero)
# ═══════════════════════════════════════════════════

s = prs.slides.add_slide(BLANK)
gradient_bg(s, BG_DARK, RGBColor(0x00, 0x18, 0x34))
accent_bar(s, 0, W, ACCENT, Pt(5))

# Subtle radial glow circle
glow = circle_text(s, Inches(4.2), Inches(-2.5), Inches(20), '', ACCENT)
# Make semi-transparent
sp = glow._element
spPr = sp.find(qn('a:spPr'))
if spPr is not None:
    sf = spPr.find(qn('a:solidFill'))
    if sf is not None:
        srgb = sf.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = parse_xml(
                '<a:alpha val="8000" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            )
            srgb.append(alpha)

# Hero text
text(s, M, Inches(1.8), Inches(11.5), Inches(1.6), '逐念而行',
     Pt(58), WHITE, bold=True)
text(s, M, Inches(3.4), Inches(11.5), Inches(0.6), '神经重连AI系统',
     Pt(26), ACCENT_LIGHT)
line(s, M, Inches(4.15), Inches(1.8), ACCENT, Pt(3))

# Tagline
text(s, M, Inches(4.6), Inches(11), Inches(0.5),
     '人机共生  ·  从大脑到肌肉的AI  ·  赛道定义者',
     Pt(15), GRAY)
text(s, M, Inches(5.25), Inches(11), Inches(0.45),
     '赵子睿  创始人 & CEO  |  2026年6月  |  种子轮',
     Pt(12), GRAY_DARK)

# Confidential notice
text(s, Inches(8.8), Inches(7.0), Inches(4.0), Inches(0.3),
     'Confidential — For Investor Use Only',
     Pt(8), GRAY_DARK, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '01', '9400万人被困在神经损伤里')

# Left: problem description
tf = tb(s, M, Inches(1.6), Inches(7.8), Inches(4.8))
para(tf, '全球9400万中风幸存者，年新发1200万。其中66%存在步行功能障碍。',
     Pt(16), DARK_TEXT, bold=True, first=True, space_after=Pt(14))
para(tf, '问题的本质：大脑还记得怎么走，但指令传不到肌肉。这不是机械问题，是神经连接问题。',
     Pt(13), GRAY_DARK, space_after=Pt(14))
para(tf, '现有方案二十年停滞不前：', Pt(14), ACCENT, bold=True, space_after=Pt(10))

bullets = [
    '外骨骼太重（12–25kg），穿戴痛苦，患者拒绝长期使用',
    '康复机器人依赖治疗师手动调参，中国每10万人仅3.6名康复治疗师',
    'SCIRE荟萃分析（2022）：现有外骨骼 vs 传统物理治疗，"证据尚不充分"',
    'Ekso FY2025收入$12.8M，SG&A占134%，连续两年下滑——硬件差异化已触天花板',
]
for b in bullets:
    para(tf, f'•  {b}', Pt(11), GRAY_DARK, space_after=Pt(6))

para(tf, '', Pt(6), space_after=Pt(8))
para(tf, '根本原因：所有竞品都在把神经损伤当成机械问题来解决。这是一个范式错误。',
     Pt(14), ACCENT, bold=True)

# Right: stat cards
stats = [
    ('9400万', '全球中风幸存者', 'GBD 2021'),
    ('66%', '步行功能障碍率', 'World Stroke Org'),
    ('3.6/10万', '康复治疗师密度', '中国 vs 高收入国家10–30'),
]
for i, (num, label, desc) in enumerate(stats):
    y = Inches(1.7) + i * Inches(1.7)
    metric_card(s, Inches(9.3), y, Inches(3.3), Inches(1.5), num, label)
    text(s, Inches(9.5), y + Inches(1.12), Inches(2.9), Inches(0.25), desc,
         Pt(8), GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 3 — SOLUTION
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '02', 'Nexum One：从大脑到肌肉的AI')
text(s, M, Inches(1.55), CW, Inches(0.35),
     '不卖硬件参数。卖的是AI对"你"的理解能力。',
     Pt(14), GRAY_DARK)

# 4-step diagram
steps = [
    ('1', '读取意图', 'EEG-Sense 8ch干电极头带(80g)\n检测运动准备电位(MRCP)\n非侵入，佩戴像运动头带\nCMU Bin He Lab 2025验证',
     ACCENT),
    ('2', '生成策略', '个性化AI引擎\n416肌肉+206骨骼数字孪生\nRL强化学习日产53亿步态\n每患者独立策略模型',
     ACCENT),
    ('3', '执行动作', 'NeuroSuit混合架构\n碳纤维锚定+线缆驱动\n<1.5kg总重, 8Nm持续\n穿戴<3分钟',
     ACCENT),
    ('4', '持续进化', '传感器反馈闭环\n策略每周自适应更新\n联邦学习隐私保护\n越多人使用AI越强',
     ACCENT),
]
for i, (num, title, desc, clr) in enumerate(steps):
    x = M + i * Inches(3.05)
    y = Inches(2.1)
    # Card
    card(s, x, y, Inches(2.85), Inches(4.5), GRAY_LIGHT, None)
    # Accent top border
    line(s, x, y, Inches(2.85), clr, Pt(3))
    # Number circle
    circle_text(s, x + Inches(0.15), y + Inches(0.18), Inches(0.45), num, clr, WHITE)
    # Title
    text(s, x + Inches(0.15), y + Inches(0.85), Inches(2.55), Inches(0.35), title,
         Pt(17), DARK_TEXT, bold=True)
    # Description
    lines = desc.split('\n')
    tf2 = tb(s, x + Inches(0.15), y + Inches(1.35), Inches(2.55), Inches(2.8))
    for j, l in enumerate(lines):
        para(tf2, l, Pt(10), GRAY_DARK, first=(j == 0), space_after=Pt(4))


# ═══════════════════════════════════════════════════
# SLIDE 4 — WHY NOW
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '03', '为什么是现在？四个条件同时就绪')

conditions = [
    ('非侵入BCI达到可用精度',
     'CMU Bin He Lab 2025 Nature Comms发表非侵入EEG解码精细运动意图。下肢步态离散意图（走/停）已被Gwin, Wagner, Sburlea等多篇独立研究验证。中国2025年3月发布非侵入BCI医疗服务定价¥960–966/次。'),
    ('可穿戴硬件工程成熟',
     '学术界已验证的混合架构（碳纤维锚定+织物包裹+线缆驱动）：力传输效率>70%，总重<1.5kg。线缆驱动临床安全性经20年验证。柔性织物和轻量化材料不再是瓶颈。'),
    ('个性化AI工程化条件成熟',
     '肌骨模型（416肌肉+206骨骼+Hill型柔顺肌腱）+ GPU并行强化学习（2048仿真世界，日产53亿步态样本，零人工标注）。"为每个患者生成专属康复策略"从学术论文变为可部署软件。'),
    ('政策框架全面到位',
     '中国脑机接口列入国家未来产业、写入政府工作报告。美国Medicare 2024年建立个人外骨骼支付标准$91,032/台。日本HAL纳入医保。德国裁定扩大外骨骼覆盖。四国已建立支付框架。'),
]
for i, (title, desc) in enumerate(conditions):
    x = M + i * Inches(3.05)
    y = Inches(1.7)
    # Card
    card(s, x, y, Inches(2.85), Inches(5.2), GRAY_LIGHT, None)
    line(s, x, y, Inches(2.85), ACCENT, Pt(3))
    # Step indicator
    circle_text(s, x + Inches(1.1), y + Inches(0.12), Inches(0.5), str(i + 1),
                ACCENT, WHITE)
    # Title
    text(s, x + Inches(0.12), y + Inches(0.8), Inches(2.6), Inches(0.55), title,
         Pt(13), DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    # Divider
    line(s, x + Inches(0.6), y + Inches(1.45), Inches(1.4), ACCENT, Pt(1))
    # Description
    tf3 = tb(s, x + Inches(0.10), y + Inches(1.6), Inches(2.65), Inches(3.3))
    para(tf3, desc, Pt(9), GRAY_DARK, first=True)

# Bottom highlight
rect(s, M, Inches(7.0), CW, Inches(0.35), HIGHLIGHT)
text(s, M + Inches(0.1), Inches(7.02), CW - Inches(0.2), Inches(0.3),
     '这四个条件——感知精度、硬件成熟度、AI工程化、政策框架——缺任何一个，人机共生都是科幻。2026年，四个条件全部就绪。',
     Pt(10), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 5 — PRODUCT
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '04', 'Nexum One：首款神经重连AI系统')

modules = [
    ('EEG-Sense 头带', '8通道干电极 · 80g\nBLE 5.0 · ADS1299芯片\n250SPS采样率 · 24-bit\n佩戴像运动头带',
     'ppt_render_2_headband.jpg'),
    ('NeuroSuit 穿戴界面', '碳纤维髋部锚定\n鲍登线缆驱动\n<1.5kg 总重\n8Nm持续 / 15Nm峰值扭矩',
     'ppt_render_3_suit.jpg'),
    ('Nexum One 完整系统', 'EEG头带 + NeuroSuit\n+ 控制盒 + App\nClinical版 $18,000\nHome版 $3,999+$499/年',
     'ppt_render_1_system.jpg'),
]
for i, (name, spec, img_name) in enumerate(modules):
    x = M + i * Inches(4.05)
    y = Inches(1.6)
    card(s, x, y, Inches(3.85), Inches(3.9), GRAY_LIGHT, ACCENT)

    # Image placeholder
    img_path = os.path.join(DOCS_DIR, img_name)
    if not try_img(s, img_path, x + Inches(0.15), y + Inches(0.15),
                   Inches(3.55), Inches(1.8)):
        rect(s, x + Inches(0.15), y + Inches(0.15), Inches(3.55), Inches(1.8),
             GRAY_CARD)
        text(s, x + Inches(0.3), y + Inches(0.85), Inches(3.25), Inches(0.35),
             name, Pt(14), GRAY, align=PP_ALIGN.CENTER)

    text(s, x + Inches(0.15), y + Inches(2.1), Inches(3.55), Inches(0.4), name,
         Pt(15), DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    line(s, x + Inches(0.8), y + Inches(2.55), Inches(2.25), ACCENT, Pt(1))
    text(s, x + Inches(0.15), y + Inches(2.7), Inches(3.55), Inches(1.1), spec,
         Pt(10), GRAY_DARK, align=PP_ALIGN.CENTER)

# Spec strip at bottom
specs = [
    ('8ch', '干电极EEG'),
    ('80g', '头带重量'),
    ('<1.5kg', '系统总重'),
    ('<185ms', '端到端延迟'),
    ('8/15Nm', '持续/峰值扭矩'),
    ('48Wh', '续航>4小时'),
]
for i, (val, label) in enumerate(specs):
    x = M + i * Inches(2.0)
    text(s, x, Inches(5.75), Inches(1.8), Inches(0.35), val,
         Pt(22), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    text(s, x, Inches(6.15), Inches(1.8), Inches(0.25), label,
         Pt(9), GRAY_DARK, align=PP_ALIGN.CENTER)

# Pricing bar
rect(s, M, Inches(6.6), CW, Inches(0.35), HIGHLIGHT)
text(s, M + Inches(0.1), Inches(6.62), CW - Inches(0.2), Inches(0.3),
     '同一套AI系统，两个版本，一条数据链：Clinical版$18K（医院）+ Home版$3,999+$499/年（家用）——数据不中断',
     Pt(11), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 6 — TECHNOLOGY
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '05', '技术架构：AI驱动的神经重连闭环')
text(s, M, Inches(1.55), CW, Inches(0.35),
     '四层架构，每层技术均已独立验证。逐念而行的贡献是把它们整合成完整的AI系统。',
     Pt(13), GRAY_DARK)

layers = [
    ('感知层', 'EEG + sEMG 双模态神经接口',
     'EEG-Sense: 8ch干电极，检测运动准备电位(MRCP)\nCMU Bin He 2025 Nature Comms验证非侵入EEG解码可行性\nsEMG: 多通道织物电极，<100ms精细力矩解码\n三级降级安全: EEG+sEMG → sEMG+IMU → 无助力被动支撑'),
    ('决策层', '个性化AI引擎',
     'DL提取EEG特征(CNN/TCN) → MSK数字孪生(416肌肉+206骨骼+Hill肌腱)\n→ 2048仿真世界并行RL训练 → 日产53亿步态 → 每患者独立策略\n→ 传统控制(PID/MPC)毫秒级安全执行\n本周策略与上周不同——因为这个患者的身体已经不同'),
    ('执行层', '混合式机械架构',
     '碳纤维髋部锚定 → 力传输效率>70%\n鲍登线缆驱动 → 电机位于腰部(重心附近)，关节处仅轻量锚点\n标准化关节电机 → 扭矩密度40–50Nm/kg，双电机+驱动器+线缆<1.5kg\n织物包裹穿戴界面 → 外观像压缩裤，解决穿戴耻感'),
    ('进化层', '联邦学习数据飞轮',
     '原始生理信号本地处理 → 仅脱敏模型改善数据上传\n联邦学习架构(pFedAC) → 隐私保护\n每多一个患者使用 → AI就更理解人 → 产品更强\n这是硬件公司无法复制的数据网络效应'),
]
for i, (label, subtitle, desc) in enumerate(layers):
    y = Inches(1.9) + i * Inches(1.28)
    # Label card
    card(s, M, y, Inches(1.85), Inches(1.15), GRAY_LIGHT, ACCENT)
    text(s, M + Inches(0.12), y + Inches(0.1), Inches(1.6), Inches(0.3), label,
         Pt(14), ACCENT, bold=True)
    text(s, M + Inches(0.12), y + Inches(0.45), Inches(1.6), Inches(0.5), subtitle,
         Pt(9), DARK_TEXT, bold=True)

    # Detail text
    detail_lines = desc.split('\n')
    tf4 = tb(s, Inches(2.95), y + Inches(0.05), Inches(9.4), Inches(1.15))
    for j, dl in enumerate(detail_lines):
        para(tf4, dl, Pt(9), GRAY_DARK, first=(j == 0), space_after=Pt(1))


# ═══════════════════════════════════════════════════
# SLIDE 7 — MARKET
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '06', '市场：TAM随每阶段指数级扩大')

phases = [
    ('PHASE 1', '2028–2030', '$6.6B', '全球神经康复',
     '9400万中风 × 66%障碍 × 30%适配 × 5%渗透 × $3,999\n+ 8000家医院 × 10%渗透 × $18K\n\n中国院内首轮TAM $43–65M\n程天已验证: 800+医院, 92万人次'),
    ('PHASE 2', '2030–2033', '$15B', '全球老年行动力',
     '中国3亿60+人口(国家统计局2025)\n× 20%肌少症 × 3%渗透 × ¥5,000\n+ 美日欧老龄化市场\n日本29% 65+, 欧盟21% 65+'),
    ('PHASE 3', '2033+', '$33.7B+', '全人群行动增强',
     '全球可穿戴外骨骼市场$33.7B\nCAGR 32% (Research & Markets)\n户外/通勤/职业防护\n人机共生新品类: 增量远超已有市场'),
]
for i, (phase_label, years, tam, title, desc) in enumerate(phases):
    x = M + i * Inches(4.05)
    bg = HIGHLIGHT if i == 2 else GRAY_LIGHT
    card(s, x, Inches(1.7), Inches(3.8), Inches(5.1), bg, ACCENT)
    # Phase label
    text(s, x + Inches(0.15), Inches(1.8), Inches(3.5), Inches(0.3),
         f'{phase_label}  {years}', Pt(10), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    # TAM
    text(s, x + Inches(0.15), Inches(2.25), Inches(3.5), Inches(0.65), tam,
         Pt(42), ACCENT, bold=True, align=PP_ALIGN.CENTER)
    line(s, x + Inches(0.7), Inches(3.0), Inches(2.4), ACCENT, Pt(1))
    # Title
    text(s, x + Inches(0.15), Inches(3.15), Inches(3.5), Inches(0.35), title,
         Pt(14), DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)
    # Description
    tf5 = tb(s, x + Inches(0.12), Inches(3.6), Inches(3.55), Inches(3.0))
    for j, dl in enumerate(desc.split('\n')):
        para(tf5, dl, Pt(9), GRAY_DARK, first=(j == 0), space_after=Pt(2))

# Bottom bar
rect(s, M, Inches(6.95), CW, Inches(0.4), GRAY_LIGHT)
text(s, M + Inches(0.1), Inches(6.97), Inches(5.5), Inches(0.3),
     '已有市场: BCI $2.4B(15% CAGR) + 康复机器人 $1.8B(22%) + 外骨骼 $33.7B(32%)',
     Pt(9), ACCENT, bold=True)
text(s, Inches(7), Inches(6.97), Inches(5.5), Inches(0.3),
     '全球支付: Medicare $91,032/台 · 中国BCI ¥960–966/次 · 日本医保 · 德国裁定',
     Pt(9), GRAY_DARK, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 8 — COMPETITION
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '07', '竞争：我们在"AI理解人"的维度上')

headers = ['竞争维度', '逐念而行 Nexum', 'Ekso / ReWalk', 'HAL (Cyberdyne)', '程天科技 UGO']
rows = [
    ['EEG运动意图解码',    '8ch干电极 ✓',          '✗',           '✗ (仅EMG)',     '多模态融合 ✓'],
    ['sEMG神经接口',        '织物电极 ✓',           '✗',           '✓',             '✓'],
    ['EEG+sEMG双模态融合',  '独家 ✓',               '✗',           '✗',             '✗'],
    ['MSK数字孪生建模',     '416肌肉+Hill肌腱 ✓',    '✗',           '✗',             '✗'],
    ['RL个性化康复策略',    '每患者独立策略 ✓',       '✗',           '✗',             '✗'],
    ['联邦学习隐私保护',    'pFedAC架构 ✓',          '✗',           '✗',             '✗'],
    ['产品价格',            '$18K / $3,999',        '$75K–100K',   '$60K–80K',      '¥10万+(医院)'],
    ['核心定位',            'AI公司',               '硬件公司',     '硬件公司',       '硬件集成'],
]
col_w = [Inches(2.3), Inches(2.6), Inches(2.15), Inches(2.15), Inches(2.3)]
add_table(s, M, Inches(1.6), Inches(11.5), Inches(4.6), headers, rows, col_w, Pt(8.5))

# Key insight card
rect(s, M, Inches(6.35), CW, Inches(0.95), HIGHLIGHT)
line(s, M, Inches(6.35), Inches(0.06), ACCENT, Pt(6))
text(s, M + Inches(0.2), Inches(6.45), CW - Inches(0.4), Inches(0.75),
     '核心差异：Ekso/ReWalk/HAL卖硬件参数——第一万台和第一台没有本质区别。逐念而行卖AI对人的理解——每多一个患者，AI就更强一分。范式竞争，不是参数竞争。',
     Pt(11), DARK_TEXT, bold=True)


# ═══════════════════════════════════════════════════
# SLIDE 9 — BUSINESS MODEL
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '08', '商业模式：硬件入口 + 订阅飞轮')

# Left column: revenue stream cards
streams = [
    ('Clinical 设备销售', '$18,000/台', 'B2B高客单价\n目标: 三甲医院康复科 + 康复专科医院\n中国4,500家目标医院\n2029年目标400台'),
    ('Home 硬件 + 订阅', '$3,999 + $499/年', 'DTC直接触达患者\n订阅收入: 持续软件更新 + AI策略\n续订率目标 >50%\n2031年8,000台×50%=$2M高毛利'),
    ('SaaS 数据服务', '$50,000+/年', '医院科研协作平台\n脱敏数据集 + AI辅助科研\n临床试验管理\n长尾高毛利收入'),
]
for i, (name, price, detail) in enumerate(streams):
    x = M
    y = Inches(1.6) + i * Inches(1.65)
    card(s, x, y, Inches(5.6), Inches(1.5), GRAY_LIGHT, ACCENT)
    # Price
    text(s, x + Inches(0.15), y + Inches(0.1), Inches(2.2), Inches(0.35), name,
         Pt(13), DARK_TEXT, bold=True)
    text(s, x + Inches(0.15), y + Inches(0.45), Inches(2.2), Inches(0.35), price,
         Pt(16), ACCENT, bold=True)
    # Detail
    detail_lines = detail.split('\n')
    tf6 = tb(s, x + Inches(2.5), y + Inches(0.1), Inches(2.9), Inches(1.3))
    para(tf6, '', Pt(2), first=True)
    for j, dl in enumerate(detail_lines):
        para(tf6, f'•  {dl}', Pt(9), GRAY_DARK, space_after=Pt(1))

# Right column: GTM + unit economics
rx = Inches(7.0)

# Customer segments
card(s, rx, Inches(1.6), Inches(5.5), Inches(2.6), GRAY_LIGHT, ACCENT)
text(s, rx + Inches(0.15), Inches(1.68), Inches(5.2), Inches(0.3), '客户分层与GTM策略',
     Pt(13), DARK_TEXT, bold=True)
segments = [
    'Tier 1: 三甲医院康复科 (中国4,500家) → 直销团队',
    'Tier 2: 康复专科医院 (中国800家) → 直销+学术推广',
    'Tier 3: 居家用户 (2029+) → 线上DTC + 医生推荐',
    '海外: 东南亚+中东代理 → 2029年启动',
]
tf7 = tb(s, rx + Inches(0.15), Inches(2.1), Inches(5.2), Inches(1.8))
for j, sg in enumerate(segments):
    para(tf7, f'•  {sg}', Pt(9), GRAY_DARK, first=(j == 0), space_after=Pt(4))

# Unit economics
card(s, rx, Inches(4.4), Inches(5.5), Inches(2.8), GRAY_LIGHT, ACCENT)
text(s, rx + Inches(0.15), Inches(4.48), Inches(5.2), Inches(0.3), '单位经济模型（Home版）',
     Pt(13), DARK_TEXT, bold=True)
unit_items = [
    ('BOM成本 (量产@100台)', '¥2,980'),
    ('BOM成本 (量产@1,000台)', '¥2,080'),
    ('首年硬件毛利', '~47% (@$3,999)'),
    ('LTV (3年订阅)', '$4,998'),
    ('CAC (线上DTC)', '~$1,200'),
    ('LTV/CAC', '~4.2x'),
]
for i, (label, val) in enumerate(unit_items):
    yy = Inches(4.9) + i * Inches(0.35)
    text(s, rx + Inches(0.15), yy, Inches(3.5), Inches(0.3), label,
         Pt(9), GRAY_DARK)
    text(s, rx + Inches(3.8), yy, Inches(1.5), Inches(0.3), val,
         Pt(9), DARK_TEXT, bold=True, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════
# SLIDE 10 — FINANCIALS
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '09', '五年财务预测')

# P&L Table
fin_headers = ['', '2027', '2028', '2029', '2030', '2031']
fin_rows = [
    ['Clinical交付(台)', '—（预临床）', '150×$16K', '400×$15K', '1,000×$14K', '2,500×$13K'],
    ['Home交付(台)', '—', '100×$4K', '500×$3.8K', '2,000×$3.5K', '8,000×$3.2K'],
    ['总收入', '$0', '$2.8M', '$7.9M', '$21.0M', '$58.1M'],
    ['毛利率', '—', '52%', '58%', '63%', '68%'],
    ['净利润', '−$2.5M', '−$2.0M', '−$0.5M', '$3.0M', '$15.0M'],
]
fin_col_w = [Inches(2.2), Inches(1.7), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.1)]
add_table(s, M, Inches(1.5), Inches(12.0), Inches(2.6), fin_headers, fin_rows,
          fin_col_w, Pt(9))

# Key metrics
metrics_data = [
    ('盈亏平衡点', '~700台/年 (2029 Q4–2030 Q1)'),
    ('Home版BOM成本', '¥2,980 (100台量产) → ¥2,080 (1,000台)'),
    ('Clinical版BOM成本', '¥3,280 (100台量产) → ¥2,280 (1,000台)'),
    ('S&M费用占比', '30% (vs 竞品40–70%，因差异化在软件而非硬件)'),
    ('Home订阅增量收入', '2031年 8,000台 × 50%续订 × $499 = +$2M (高毛利)'),
    ('研发投入占比', '2027: 80% → 2029: 35% → 2031: 15%'),
]
for i, (label, val) in enumerate(metrics_data):
    row = i // 2
    col = i % 2
    mx = M + col * Inches(6.2)
    my = Inches(4.4) + row * Inches(0.6)
    card(s, mx, my, Inches(5.9), Inches(0.5), GRAY_LIGHT, None)
    text(s, mx + Inches(0.1), my + Inches(0.08), Inches(1.8), Inches(0.3), label,
         Pt(9), DARK_TEXT, bold=True)
    text(s, mx + Inches(1.9), my + Inches(0.08), Inches(3.8), Inches(0.3), val,
         Pt(9), GRAY_DARK)

# Breakeven highlight
rect(s, M, Inches(5.75), CW, Inches(0.35), HIGHLIGHT)
text(s, M + Inches(0.1), Inches(5.77), CW - Inches(0.2), Inches(0.3),
     '盈亏平衡预计在2029 Q4–2030 Q1达到，年销~700台。净利2027年转正前共需约$5M融资。',
     Pt(10), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 11 — TEAM
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '10', '团队：六人核心')
text(s, M, Inches(1.55), CW, Inches(0.35),
     '优先级：临床决定NMPA · EEG决定量产 · 注册决定审评速度 · 设计决定用户留存',
     Pt(12), GRAY_DARK)

members = [
    ('赵子睿', '创始人 & CEO', 'Columbia EE MSc (ML)\n华为美国研究院Futurewei IoT Lab\n欣网/智擎CTO 8年\nFAI患者——产品原动力', True),
    ('汪牧星', '算法 · 联邦学习', 'NEU PhD Candidate\nEdinburgh MSc Distinction\npFedAC第一作者', True),
    ('临床负责人', '招募中 · 最高优先级', '康复医学主任医师\n临床试验设计\nNMPA方案撰写经验', False),
    ('EEG硬件负责人', '招募中', '脑电采集系统设计\nTI ADS1299经验\n干电极优化+产品化', False),
    ('NMPA注册负责人', '招募中', '二类创新器械注册\nISO 13485体系\n创新器械申报经验', False),
    ('产品设计师', '招募中', '可穿戴设备工业设计\n用户体验研究\n穿戴体验=留存关键', False),
]
for i, (name, role, bio, active) in enumerate(members):
    x = M + i * Inches(2.05)
    y = Inches(1.85)
    bg = WHITE if active else GRAY_LIGHT
    bdr = ACCENT if active else None
    card(s, x, y, Inches(1.85), Inches(3.5), bg, bdr)
    # Avatar circle
    circle_text(s, x + Inches(0.6), y + Inches(0.1), Inches(0.55),
                name[0], ACCENT if active else GRAY, WHITE)
    # Name
    text(s, x + Inches(0.08), y + Inches(0.75), Inches(1.7), Inches(0.35), name,
         Pt(13), DARK_TEXT if active else GRAY, bold=True, align=PP_ALIGN.CENTER)
    # Role
    text(s, x + Inches(0.08), y + Inches(1.15), Inches(1.7), Inches(0.4), role,
         Pt(8), ACCENT if active else GRAY, align=PP_ALIGN.CENTER)
    # Bio
    bio_lines = bio.split('\n')
    tf8 = tb(s, x + Inches(0.08), y + Inches(1.65), Inches(1.7), Inches(1.6))
    for j, bl in enumerate(bio_lines):
        para(tf8, bl, Pt(7), GRAY_DARK if active else GRAY,
             first=(j == 0), space_after=Pt(2))

# Academic partners
rect(s, M, Inches(5.6), CW, Inches(1.6), GRAY_LIGHT)
text(s, M + Inches(0.15), Inches(5.68), CW - Inches(0.3), Inches(0.3),
     '学术合作网络', Pt(13), DARK_TEXT, bold=True)
text(s, M + Inches(0.15), Inches(6.05), CW - Inches(0.3), Inches(0.3),
     'CMU Bin He Lab (非侵入EEG运动解码) · NEU NeuMove Lab (肌骨仿真) · 清华/浙大/西湖大学 (BCI与神经工程)',
     Pt(9), GRAY_DARK)
text(s, M + Inches(0.15), Inches(6.4), CW - Inches(0.3), Inches(0.3),
     '杨朋昆 (清华统计系副教授, pFedAC共同作者) · 苏丽丽 (NEU助理教授, NSF CAREER, 联邦学习)',
     Pt(9), GRAY_DARK)


# ═══════════════════════════════════════════════════
# SLIDE 12 — FUNDING
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '11', '融资路线：每轮对应可验证里程碑')

rounds = [
    ('种子轮', '2026 Q3', '¥800万',
     '原型2台(¥200万)\n团队8人12月(¥300万)\n预临床启动(¥200万)\n运营储备(¥100万)',
     '可演示原型\n3家医院预临床协议\nNMPA分类界定申请', False),
    ('天使轮', '2027 Q1', '¥2,500万',
     '10台样机(¥500万)\n团队12人18月(¥800万)\n30例预临床(¥600万)\nNMPA申报(¥300万)',
     '预临床数据发表\nNMPA创新器械正式受理', False),
    ('Pre-A', '2028 Q1', '$7M (≈¥5,000万)',
     '量产备货(¥2,000万)\n50家医院渠道(¥1,500万)\nHome版发布(¥1,000万)\n团队30人(¥500万)',
     'NMPA获批\n首批商业化交付\n年收入$2.8M', False),
    ('A轮', '2029 H2', '$20M',
     'FDA 510(k)+CE申报\n海外渠道建设\n规模化生产\n团队50人',
     '海外准入\n年销过千台\n盈亏平衡', True),
]
for i, (name, time, amount, use, milestone, is_last) in enumerate(rounds):
    x = M + i * Inches(3.05)
    bg = HIGHLIGHT if is_last else GRAY_LIGHT
    card(s, x, Inches(1.7), Inches(2.85), Inches(5.2), bg, ACCENT)

    # Round name
    text(s, x + Inches(0.15), Inches(1.8), Inches(2.55), Inches(0.35), name,
         Pt(18), DARK_TEXT, bold=True)
    text(s, x + Inches(0.15), Inches(2.2), Inches(2.55), Inches(0.25), time,
         Pt(10), GRAY_DARK)
    text(s, x + Inches(0.15), Inches(2.5), Inches(2.55), Inches(0.45), amount,
         Pt(26), ACCENT, bold=True)
    line(s, x + Inches(0.3), Inches(3.0), Inches(2.25), ACCENT, Pt(1))

    # Use of funds
    text(s, x + Inches(0.15), Inches(3.15), Inches(2.55), Inches(0.25), '用途',
         Pt(9), DARK_TEXT, bold=True)
    for j, line_text in enumerate(use.split('\n')):
        text(s, x + Inches(0.15), Inches(3.45 + j * 0.3), Inches(2.55), Inches(0.25),
             line_text, Pt(8), GRAY_DARK)

    # Milestone / output
    my_start = Inches(4.85)
    text(s, x + Inches(0.15), my_start, Inches(2.55), Inches(0.25), '产出里程碑',
         Pt(9), ACCENT, bold=True)
    rect(s, x + Inches(0.1), my_start + Inches(0.28), Inches(2.7), Inches(0.8),
         WHITE)
    tf9 = tb(s, x + Inches(0.2), my_start + Inches(0.3), Inches(2.5), Inches(0.75))
    for j, ml in enumerate(milestone.split('\n')):
        para(tf9, f'•  {ml}', Pt(8), DARK_TEXT if is_last else GRAY_DARK,
             first=(j == 0), space_after=Pt(2))

# Bottom summary
rect(s, M, Inches(7.0), CW, Inches(0.35), HIGHLIGHT)
text(s, M + Inches(0.1), Inches(7.02), CW - Inches(0.2), Inches(0.3),
     '当前：种子轮¥800万寻找领投方。12个月可验证的技术里程碑 → 原型可演示 → 预临床数据 → 为天使轮¥2,500万建立定价基础。',
     Pt(10), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 13 — ROADMAP
# ═══════════════════════════════════════════════════

s = make_content_slide(prs, '12', '路线图：2026–2031')

# Timeline horizontal line
line(s, M, Inches(3.5), CW, ACCENT, Pt(3))

# Phase markers and milestones
phases_roadmap = [
    ('Phase 1', '2026–2027', '研发与验证',
     ['种子轮关闭, 团队组建', '功能原型完成', '预临床启动 (3家医院)', 'NMPA创新器械认定', '30例预临床数据发表'],
     Inches(1.7), ACCENT),
    ('Phase 2', '2028–2029', '获批与商业化',
     ['注册检验完成', 'NMPA获批 (二类创新)', '首批Clinical交付150台', 'Home版正式发布', '盈亏平衡 (~700台/年)'],
     Inches(4.0), ACCENT),
    ('Phase 3', '2030–2031', '规模化与国际化',
     ['年销10,000+台', 'FDA 510(k) 提交', 'CE认证启动', '东南亚+中东渠道', '启动全人群行动增强市场'],
     Inches(7.0), ACCENT),
]
for i, (phase_label, years, phase_title, milestones, cx, clr) in enumerate(phases_roadmap):
    # Dot on timeline
    circle_text(s, cx, Inches(3.35), Inches(0.3), '', clr, clr)

    # Header above timeline
    text(s, cx - Inches(0.8), Inches(1.7), Inches(2.1), Inches(0.25), phase_label,
         Pt(10), clr, bold=True, align=PP_ALIGN.CENTER)
    text(s, cx - Inches(0.8), Inches(1.95), Inches(2.1), Inches(0.25), years,
         Pt(9), GRAY_DARK, align=PP_ALIGN.CENTER)
    text(s, cx - Inches(0.8), Inches(2.2), Inches(2.1), Inches(0.3), phase_title,
         Pt(12), DARK_TEXT, bold=True, align=PP_ALIGN.CENTER)

    # Milestones below timeline
    card(s, cx - Inches(0.95), Inches(3.9), Inches(2.1), Inches(2.8), GRAY_LIGHT, clr)
    tf10 = tb(s, cx - Inches(0.85), Inches(3.95), Inches(1.9), Inches(2.7))
    para(tf10, '', Pt(2), first=True)
    for ml in milestones:
        para(tf10, f'•  {ml}', Pt(9), GRAY_DARK, space_after=Pt(6))

# Bottom note
rect(s, M, Inches(6.95), CW, Inches(0.35), HIGHLIGHT)
text(s, M + Inches(0.1), Inches(6.97), CW - Inches(0.2), Inches(0.3),
     '每18个月一个里程碑审查节点，融资节奏与产品节奏对齐，减少不必要的稀释。',
     Pt(10), ACCENT, bold=True, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SLIDE 14 — VISION (Dark Closing)
# ═══════════════════════════════════════════════════

s = prs.slides.add_slide(BLANK)
gradient_bg(s, BG_DARK, RGBColor(0x00, 0x20, 0x50))
accent_bar(s, 0, W, ACCENT, Pt(5))

# Vision statement
text(s, M, Inches(1.0), Inches(11.5), Inches(1.8),
     '2026年，AI从数字世界进入了物理世界。\n逐念而行专注于其中尚未被充分探索的方向——\n让AI理解人的身体，而不只是造AI的机器身体。',
     Pt(24), WHITE, align=PP_ALIGN.CENTER)

line(s, Inches(5.5), Inches(3.1), Inches(2.3), ACCENT, Pt(2))

# Vision timeline
vision_items = [
    ('2026–2028', '定义品类',
     '100家医院 · 30例预临床 → 200例上市后研究\n临床论文 → 进入康复指南\n让"人机共生"出现在顶级期刊和政策文件中'),
    ('2029–2030', '品类扩展',
     '医院→家庭闭环 · 从神经重连到通用意图驱动助力\n中国3亿60+老人，老年行动力市场\n启动FDA + CE，海外准入'),
    ('2031+', '成为基础设施',
     '百万用户 · 从修复到增强\n户外运动 · 都市通勤 · 职业防护\n人机共生像智能手机一样成为日常生活一部分'),
]
for i, (year, title, desc) in enumerate(vision_items):
    y = Inches(3.5) + i * Inches(1.1)
    # Year
    rect(s, M, y, Inches(1.4), Inches(0.4), ACCENT)
    text(s, M + Inches(0.05), y + Inches(0.03), Inches(1.3), Inches(0.35), year,
         Pt(12), WHITE, bold=True, align=PP_ALIGN.CENTER)
    # Title
    text(s, M + Inches(1.7), y, Inches(2.0), Inches(0.4), title,
         Pt(16), ACCENT_LIGHT, bold=True)
    # Description
    text(s, M + Inches(4.0), y - Inches(0.05), Inches(8.0), Inches(0.95), desc,
         Pt(9), GRAY)

# Contact card
y_card = Inches(6.2)
rect(s, Inches(2.5), y_card, Inches(8.3), Inches(1.0), CARD_DARK)
text(s, Inches(2.7), y_card + Inches(0.08), Inches(7.9), Inches(0.35),
     '逐念而行 · Nexum  |  人机共生赛道定义者',
     Pt(18), WHITE, bold=True, align=PP_ALIGN.CENTER)
text(s, Inches(2.7), y_card + Inches(0.52), Inches(7.9), Inches(0.3),
     '赵子睿  ·  zirui@nexum.ai  ·  Columbia EE MSc  ·  种子轮 ¥800万',
     Pt(11), ACCENT_LIGHT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════

os.makedirs(DOCS_DIR, exist_ok=True)
prs.save(OUTPUT)
print(f'Nexum Pitch Deck saved: {OUTPUT}')
print(f'Size: {os.path.getsize(OUTPUT) / 1024:.0f} KB  |  Slides: {len(prs.slides)}')
